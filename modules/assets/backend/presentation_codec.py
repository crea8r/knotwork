from __future__ import annotations

import json
import re
from io import BytesIO
from pathlib import Path
from typing import Any


PRESENTATION_FILE_TYPE = "presentation"
PRESENTATION_EXTS = {".ppt", ".pptx"}
PRESENTATION_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def presentation_title_from_filename(filename: str) -> str:
    stem = Path(filename).stem.replace("-", " ").replace("_", " ").strip()
    return stem.title() or "Untitled Presentation"


def default_presentation_document(title: str = "Untitled Presentation") -> dict[str, Any]:
    return {
        "kind": "presentation",
        "version": 1,
        "title": title,
        "theme": {
            "background": "#f5efe2",
            "accent": "#1f4f8c",
            "text": "#1d2733",
        },
        "slides": [
            {
                "id": "slide-1",
                "title": title,
                "subtitle": "Add a subtitle",
                "body": "",
                "bullets": [],
                "notes": "",
                "layout": "title",
            }
        ],
    }


def presentation_to_storage_content(document: dict[str, Any]) -> str:
    normalized = normalize_presentation_document(document)
    return json.dumps(normalized, indent=2)


def presentation_from_storage_content(content: str, fallback_title: str | None = None) -> dict[str, Any]:
    title = fallback_title or "Untitled Presentation"
    if not content.strip():
        return default_presentation_document(title)
    try:
        payload = json.loads(content)
    except json.JSONDecodeError:
        payload = default_presentation_document(title)
        payload["slides"][0]["body"] = content.strip()
        return payload
    return normalize_presentation_document(payload, fallback_title=title)


def normalize_presentation_document(
    payload: dict[str, Any] | None,
    fallback_title: str = "Untitled Presentation",
) -> dict[str, Any]:
    source = payload if isinstance(payload, dict) else {}
    title = str(source.get("title") or fallback_title).strip() or fallback_title
    theme_source = source.get("theme") if isinstance(source.get("theme"), dict) else {}
    theme = {
        "background": _color(theme_source.get("background"), "#f5efe2"),
        "accent": _color(theme_source.get("accent"), "#1f4f8c"),
        "text": _color(theme_source.get("text"), "#1d2733"),
    }
    raw_slides = source.get("slides")
    slides_source = raw_slides if isinstance(raw_slides, list) else []
    slides = [_normalize_slide(slide, index, title) for index, slide in enumerate(slides_source)]
    if not slides:
        slides = default_presentation_document(title)["slides"]
    return {
        "kind": "presentation",
        "version": 1,
        "title": title,
        "theme": theme,
        "slides": slides,
    }


def import_presentation_bytes(filename: str, content: bytes) -> tuple[dict[str, Any], str]:
    suffix = Path(filename).suffix.lower()
    title = presentation_title_from_filename(filename)
    if suffix == ".pptx":
        return _import_pptx(content, title), "pptx"
    if suffix == ".ppt":
        return _import_legacy_ppt(content, title), "ppt"
    raise ValueError(f"Unsupported presentation type: {suffix}")


def export_presentation_bytes(document: dict[str, Any]) -> bytes:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("python-pptx is not installed") from exc

    deck = normalize_presentation_document(document)
    prs = Presentation()
    accent = _rgb(deck["theme"]["accent"])
    text_color = _rgb(deck["theme"]["text"])

    while prs.slides:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    for slide_doc in deck["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(8.5), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_p = title_frame.paragraphs[0]
        title_run = title_p.add_run()
        title_run.text = slide_doc["title"]
        title_run.font.size = Pt(28)
        title_run.font.bold = True
        title_run.font.color.rgb = accent

        if slide_doc["subtitle"]:
            subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(8.2), Inches(0.7))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_run = subtitle_p.add_run()
            subtitle_run.text = slide_doc["subtitle"]
            subtitle_run.font.size = Pt(16)
            subtitle_run.font.color.rgb = text_color

        body_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(8.0), Inches(4.4))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.clear()
        wrote_body = False
        if slide_doc["body"]:
            for index, line in enumerate(_split_lines(slide_doc["body"])):
                paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
                paragraph.alignment = PP_ALIGN.LEFT
                run = paragraph.add_run()
                run.text = line
                run.font.size = Pt(20)
                run.font.color.rgb = text_color
                wrote_body = True
        for bullet in slide_doc["bullets"]:
            paragraph = body_frame.add_paragraph() if wrote_body else body_frame.paragraphs[0]
            wrote_body = True
            paragraph.level = 0
            run = paragraph.add_run()
            run.text = f"• {bullet}"
            run.font.size = Pt(20)
            run.font.color.rgb = text_color

        if slide_doc["notes"]:
            notes_text_frame = slide.notes_slide.notes_text_frame
            notes_text_frame.text = slide_doc["notes"]

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def presentation_summary(document: dict[str, Any]) -> str:
    deck = normalize_presentation_document(document)
    lines = [f"# {deck['title']}", ""]
    for index, slide in enumerate(deck["slides"], start=1):
        lines.append(f"## Slide {index}: {slide['title'] or 'Untitled'}")
        if slide["subtitle"]:
            lines.append(slide["subtitle"])
        if slide["body"]:
            lines.extend(_split_lines(slide["body"]))
        for bullet in slide["bullets"]:
            lines.append(f"- {bullet}")
        lines.append("")
    return "\n".join(lines).strip()


def _import_pptx(content: bytes, title: str) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is not installed") from exc

    prs = Presentation(BytesIO(content))
    slides: list[dict[str, Any]] = []
    for index, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        subtitle = ""
        body_lines: list[str] = []
        bullets: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = "\n".join(
                paragraph.text.strip()
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text and paragraph.text.strip()
            ).strip()
            if not text:
                continue
            if getattr(shape, "is_placeholder", False):
                placeholder_name = str(getattr(shape, "name", "")).lower()
                if "title" in placeholder_name and not slide_title:
                    slide_title = text.splitlines()[0].strip()
                    continue
                if "subtitle" in placeholder_name and not subtitle:
                    subtitle = text
                    continue
            for paragraph in shape.text_frame.paragraphs:
                text_value = paragraph.text.strip()
                if not text_value:
                    continue
                if getattr(paragraph, "level", 0) > 0 or text_value.startswith("• "):
                    bullets.append(text_value.removeprefix("• ").strip())
                else:
                    body_lines.append(text_value)
        if not slide_title and body_lines:
            slide_title = body_lines.pop(0)
        if not subtitle and body_lines and len(body_lines[0]) < 120:
            subtitle = body_lines.pop(0)
        slides.append(
            _normalize_slide(
                {
                    "id": f"slide-{index}",
                    "title": slide_title or f"Slide {index}",
                    "subtitle": subtitle,
                    "body": "\n".join(body_lines[:6]).strip(),
                    "bullets": bullets[:8],
                    "notes": "",
                    "layout": "title-body",
                },
                index - 1,
                title,
            )
        )
    return normalize_presentation_document({"title": title, "slides": slides}, fallback_title=title)


def _import_legacy_ppt(content: bytes, title: str) -> dict[str, Any]:
    decoded = content.decode("latin-1", errors="ignore")
    runs = re.findall(r"[A-Za-z0-9][^\x00-\x08\x0b\x0c\x0e-\x1f]{4,}", decoded)
    cleaned = [re.sub(r"\s+", " ", run).strip() for run in runs]
    cleaned = [run for run in cleaned if len(run) > 4]
    if not cleaned:
        return default_presentation_document(title)

    chunk_size = 8
    slides = []
    for index in range(0, min(len(cleaned), 40), chunk_size):
        chunk = cleaned[index:index + chunk_size]
        slides.append(
            {
                "id": f"slide-{len(slides) + 1}",
                "title": chunk[0][:80],
                "subtitle": "",
                "body": "",
                "bullets": chunk[1:6],
                "notes": "Imported from legacy .ppt with best-effort text extraction.",
                "layout": "title-body",
            }
        )
    return normalize_presentation_document({"title": title, "slides": slides}, fallback_title=title)


def _normalize_slide(slide: Any, index: int, fallback_title: str) -> dict[str, Any]:
    source = slide if isinstance(slide, dict) else {}
    title = str(source.get("title") or f"Slide {index + 1}").strip() or f"Slide {index + 1}"
    subtitle = str(source.get("subtitle") or "").strip()
    body = str(source.get("body") or "").strip()
    raw_bullets = source.get("bullets")
    bullets = [str(item).strip() for item in raw_bullets] if isinstance(raw_bullets, list) else []
    bullets = [item for item in bullets if item]
    return {
        "id": str(source.get("id") or f"slide-{index + 1}"),
        "title": title,
        "subtitle": subtitle,
        "body": body,
        "bullets": bullets,
        "notes": str(source.get("notes") or "").strip(),
        "layout": "title" if str(source.get("layout") or "") == "title" else "title-body",
    }


def _color(value: Any, fallback: str) -> str:
    if isinstance(value, str) and re.fullmatch(r"#[0-9a-fA-F]{6}", value):
        return value.lower()
    return fallback


def _rgb(value: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(value.removeprefix("#").upper())


def _split_lines(value: str) -> list[str]:
    return [line.strip() for line in value.splitlines() if line.strip()]
